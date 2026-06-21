from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


OUT = Path("/Users/liwenbin930/Desktop/bytesavor-backend/期末展演/ByteSavor_期末答辩_最终5页增强架构版.pptx")
W, H = 13.333, 7.5

COLORS = {
    "bg": RGBColor(247, 251, 248),
    "deep": RGBColor(23, 59, 46),
    "green": RGBColor(27, 127, 95),
    "mint": RGBColor(85, 201, 155),
    "amber": RGBColor(244, 185, 87),
    "purple": RGBColor(92, 86, 177),
    "muted": RGBColor(92, 107, 99),
    "pale_green": RGBColor(234, 247, 239),
    "pale_amber": RGBColor(255, 248, 234),
    "pale_purple": RGBColor(244, 243, 255),
    "white": RGBColor(255, 255, 255),
    "red": RGBColor(207, 83, 72),
}


def add_text(slide, text, x, y, w, h, size=18, color="deep", bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "PingFang SC"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = COLORS[color] if isinstance(color, str) else color
    if align:
        p.alignment = align
    return box


def add_round(slide, x, y, w, h, fill="white", line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill] if isinstance(fill, str) else fill
    if line:
        shape.line.color.rgb = COLORS[line] if isinstance(line, str) else line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_title(slide, title, subtitle):
    add_text(slide, title, 0.65, 0.32, 12, 0.5, 25, "deep", True)
    add_text(slide, subtitle, 0.67, 0.82, 11.8, 0.32, 13.2, "muted")


def add_arrow(slide, x1, y1, x2, y2, color="muted", width=2.1):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = COLORS[color] if isinstance(color, str) else color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def set_bg(slide, color="bg"):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS[color] if isinstance(color, str) else color


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, "bg")
    # Decorative editable blocks.
    add_round(slide, 9.85, -0.22, 3.3, 1.75, "pale_green")
    add_round(slide, 10.92, 5.72, 2.25, 1.28, "pale_amber")
    add_round(slide, -0.48, 5.58, 2.9, 1.7, "pale_purple")
    add_text(slide, "ByteSavor", 0.75, 0.7, 3.8, 0.5, 30, "green", True)
    add_text(slide, "基于多模态 Agent 的\n全场景饮食全链路解析系统", 0.75, 1.45, 7.25, 1.25, 34, "deep", True)
    add_text(slide, "从看见食物，到推荐、执行、记录和长期偏好学习", 0.78, 2.72, 6.5, 0.35, 15, "muted")

    # Phone-like UI cluster, all editable shapes.
    add_round(slide, 8.1, 1.0, 3.35, 5.6, "deep")
    add_round(slide, 8.33, 1.25, 2.88, 5.1, "white")
    add_text(slide, "今日饮食总控台", 8.58, 1.55, 2.0, 0.22, 12, "deep", True)
    add_round(slide, 8.58, 1.95, 2.35, 0.88, "pale_green")
    add_text(slide, "健康分 82", 8.8, 2.18, 1.2, 0.2, 14, "green", True)
    add_text(slide, "蛋白质还缺 24g", 8.82, 2.48, 1.3, 0.18, 9.5, "muted")
    for i, (label, fill, x, y) in enumerate([
        ("拍照识别", "pale_green", 8.58, 3.15),
        ("营养分析", "pale_amber", 9.78, 3.15),
        ("品质鉴定", "pale_purple", 8.58, 4.05),
        ("探店向导", "pale_green", 9.78, 4.05),
    ]):
        add_round(slide, x, y, 1.02, 0.65, fill)
        add_text(slide, label, x + 0.08, y + 0.25, 0.86, 0.15, 8.3, "deep", True, PP_ALIGN.CENTER)
    add_round(slide, 8.58, 5.05, 2.35, 0.72, "pale_amber")
    add_text(slide, "Agent：牛肉南瓜，30分钟减脂餐", 8.75, 5.29, 2.05, 0.16, 8.5, "deep", True, PP_ALIGN.CENTER)

    # B-Y-T-E strip.
    steps = [("B", "感知"), ("Y", "决策"), ("T", "执行"), ("E", "反馈")]
    for i, (letter, word) in enumerate(steps):
        x = 0.85 + i * 1.45
        add_round(slide, x, 5.6, 1.05, 0.72, "white")
        add_text(slide, letter, x + 0.18, 5.72, 0.28, 0.25, 18, "green", True, PP_ALIGN.CENTER)
        add_text(slide, word, x + 0.48, 5.78, 0.42, 0.18, 10.5, "deep", True)
        if i < len(steps) - 1:
            add_arrow(slide, x + 1.08, 5.96, x + 1.35, 5.96, "green", 1.6)
    add_round(slide, 0.78, 4.12, 6.5, 1.03, "white")
    add_text(slide, "完整闭环", 1.05, 4.34, 0.95, 0.16, 11.5, "green", True)
    add_text(slide, "拍照识别 -> 智能决策 -> 菜谱推荐 -> 清单执行 -> 确认摄入 -> 营养记录 -> 偏好学习", 1.95, 4.32, 4.8, 0.18, 10.5, "deep", True)
    add_text(slide, "这条链路保证“看见食物”以后能继续落到行动和长期记忆，而不是停在一次 API 返回。", 1.05, 4.68, 5.55, 0.18, 9.5, "muted")
    add_text(slide, "功能概览 · 技术栈/技术难点 · 软件规模 · 角色化测试 · 现场演示", 0.85, 6.62, 8.5, 0.25, 12, "muted")


def slide_roles_baseline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, "bg")
    add_title(slide, "功能概览：按角色组织，而不是堆按钮", "同一套多模态能力，覆盖家庭做饭、健康管理、买菜、探店和长期个性化")

    roles = [
        ("家庭做饭", "拍照识别食材\n校正后推荐下一餐", "pale_green"),
        ("健康管理", "个性化目标\n今日/本周营养缺口", "pale_amber"),
        ("买菜用户", "品质鉴定\n外观依据和挑选建议", "pale_purple"),
        ("探店用户", "菜品识别\n历史故事、口味技法", "pale_green"),
        ("备餐采购", "多菜谱合并\n购物清单与数量保留", "pale_amber"),
        ("长期用户", "评分原因\n偏好记忆影响推荐", "pale_purple"),
    ]
    for i, (title, body, fill) in enumerate(roles):
        x = 0.7 + (i % 3) * 2.2
        y = 1.55 + (i // 3) * 1.18
        add_round(slide, x, y, 1.95, 0.92, fill)
        add_text(slide, title, x + 0.16, y + 0.15, 1.55, 0.2, 12.5, "deep", True)
        add_text(slide, body, x + 0.16, y + 0.42, 1.6, 0.36, 9.2, "muted")

    add_text(slide, "Baseline 对比", 7.2, 1.5, 2.1, 0.3, 17, "deep", True)
    add_text(slide, "传统饮食软件", 7.2, 2.0, 1.6, 0.2, 11.5, "muted", True)
    for i, label in enumerate(["热量记录", "菜谱搜索", "购物清单", "图片识别"]):
        x = 7.2 + (i % 2) * 1.42
        y = 2.35 + (i // 2) * 0.82
        add_round(slide, x, y, 1.16, 0.55, "white", "pale_green")
        add_text(slide, label, x + 0.1, y + 0.2, 0.95, 0.12, 8.8, "muted", True, PP_ALIGN.CENTER)
    add_text(slide, "能力割裂", 8.04, 3.96, 1.0, 0.18, 10.5, "red", True, PP_ALIGN.CENTER)
    add_arrow(slide, 9.1, 3.2, 9.85, 3.2, "green", 2)

    add_text(slide, "ByteSavor", 10.05, 2.0, 1.5, 0.2, 12, "green", True)
    cx, cy = 10.78, 3.22
    core = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 0.62), Inches(cy - 0.62), Inches(1.24), Inches(1.24))
    core.fill.solid(); core.fill.fore_color.rgb = COLORS["green"]; core.line.fill.background()
    add_text(slide, "B-Y-T-E", cx - 0.46, cy - 0.11, 0.92, 0.2, 13.5, "white", True, PP_ALIGN.CENTER)
    for text, x, y, c in [("感知", cx - .32, cy - 1.0, "green"), ("决策", cx + .82, cy - .05, "amber"), ("执行", cx - .32, cy + .9, "purple"), ("反馈", cx - 1.42, cy - .05, "green")]:
        add_round(slide, x, y, 0.72, 0.3, "pale_green" if c == "green" else "pale_amber" if c == "amber" else "pale_purple")
        add_text(slide, text, x + 0.07, y + 0.09, 0.58, 0.1, 8.5, c, True, PP_ALIGN.CENTER)
    add_text(slide, "从识别到执行，再到长期记忆", 9.75, 4.58, 2.4, 0.22, 11, "muted", False, PP_ALIGN.CENTER)

    add_round(slide, 0.8, 5.35, 11.8, 0.72, "deep")
    add_text(slide, "一句话：竞品更像互相孤立的工具集合，ByteSavor 把真实饮食任务接成一条可追踪、可验证的 Agent 流水线。", 1.15, 5.58, 11.1, 0.2, 12.5, "white", True)


def slide_baseline_to_closure(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, "bg")
    add_text(slide, "BASELINE TO CLOSURE", 0.7, 0.34, 4.0, 0.25, 12, "green", True)
    add_text(slide, "从离散工具到 BYTE 闭环", 0.7, 0.67, 6.8, 0.45, 26, "deep", True)
    add_text(slide, "原始 Baseline 的核心问题不是“没有 App”，而是识别、营养、菜谱、购物、反馈互相断开。ByteSavor 的改进点是把这些节点连接成一个可执行流程。", 0.72, 1.18, 8.8, 0.45, 12.2, "muted")
    add_text(slide, "02/05", 11.55, 0.55, 1.0, 0.26, 12, "muted", True, PP_ALIGN.RIGHT)

    add_round(slide, 0.72, 1.95, 5.35, 4.35, "white")
    add_text(slide, "Baseline：离散与断层", 1.05, 2.28, 2.4, 0.26, 16, "deep", True)
    baseline = [
        ("1", "识别", "只能知道“是什么”"),
        ("2", "菜谱", "单点搜索，缺少约束"),
        ("3", "营养", "另开工具手动换算"),
        ("4", "清单", "重复整理、难导出"),
        ("5", "反馈", "评价不回流推荐"),
    ]
    for i, (num, title, body) in enumerate(baseline):
        y = 2.78 + i * 0.55
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.04), Inches(y), Inches(0.34), Inches(0.34))
        circ.fill.solid(); circ.fill.fore_color.rgb = COLORS["pale_amber"]; circ.line.fill.background()
        add_text(slide, num, 1.12, y + 0.09, 0.18, 0.1, 8.2, "amber", True, PP_ALIGN.CENTER)
        add_text(slide, title, 1.55, y + 0.06, 0.78, 0.13, 10.8, "deep", True)
        add_text(slide, body, 2.3, y + 0.06, 2.5, 0.13, 10.2, "muted")
    add_round(slide, 1.05, 5.72, 4.55, 0.38, "pale_amber")
    add_text(slide, "结果：信息到行动之间存在多次手动跳转，现场演示容易变成“打开多个工具”。", 1.22, 5.84, 4.25, 0.1, 8.7, "muted")

    add_round(slide, 6.65, 1.95, 5.95, 4.35, "deep")
    add_text(slide, "ByteSavor：Agent 全链路闭环", 7.0, 2.28, 3.8, 0.26, 16, "white", True)
    steps = [
        ("B", "Sense", "图片/文本输入", "green"),
        ("Y", "Decision", "推荐 + 营养缺口", "amber"),
        ("T", "Task", "清单合并 + 导出", "purple"),
        ("E", "Feedback", "评分学习偏好", "green"),
    ]
    for i, (letter, en, body, c) in enumerate(steps):
        x = 7.0 + i * 1.28
        add_round(slide, x, 3.05, 1.05, 1.1, "white")
        add_text(slide, letter, x + 0.25, 3.22, 0.5, 0.24, 19, c, True, PP_ALIGN.CENTER)
        add_text(slide, en, x + 0.13, 3.58, 0.8, 0.13, 8.7, "deep", True, PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.1, 3.84, 0.85, 0.22, 7.6, "muted", False, PP_ALIGN.CENTER)
        if i < 3:
            add_arrow(slide, x + 1.08, 3.6, x + 1.23, 3.6, "pale_green", 1.5)
    add_arrow(slide, 10.95, 4.25, 7.25, 4.85, "pale_green", 1.8)
    add_round(slide, 7.0, 5.12, 5.05, 0.62, "pale_green")
    add_text(slide, "每次调用保留 trace_id / stages / latency / degraded，能解释、能兜底、能复盘。", 7.28, 5.34, 4.55, 0.14, 9.2, "deep", True, PP_ALIGN.CENTER)


def slide_arch(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, "deep")
    add_text(slide, "技术框架：从前端体验到 Agent 编排", 0.65, 0.32, 12, 0.5, 25, "white", True)
    add_text(slide, "把用户看到的角色页面，映射到后端 API、Agent 编排、工具服务、数据记忆和外部模型", 0.67, 0.82, 11.8, 0.32, 13.2, RGBColor(201, 225, 213))

    layers = [
        ("前端体验层\nExperience", "首页总控台 / 拍照识别 / 营养看板 / 品质鉴定 / 探店向导 / 文本导入", "pale_green"),
        ("API 接入层\nFastAPI Routers", "auth / sense / decision / meals / feedback / nutrition / quality / guide / agent", "white"),
        ("Agent 编排层\nLangGraph Runtime", "Planner -> Tool Start -> Tool Result -> Final；输出 trace_id / stages / latency / degraded", "pale_amber"),
        ("工具服务层\nB-Y-T-E Services", "B Better Perception · Y Yielding Decisions · T Task Automation · E Evolving Feedback", "white"),
        ("数据与记忆层\nData & Memory", "MySQL: User/Profile/MealRecord/PreferenceMemory/Recipe；Redis 缓存；2576+ 菜谱数据", "pale_purple"),
        ("外部模型层\nModel Providers", "VLM 负责看图；LLM 负责意图、讲解和偏好解析；失败时记录降级而非伪造 mock", "white"),
    ]
    for i, (title, body, fill) in enumerate(layers):
        y = 1.35 + i * 0.75
        add_round(slide, 0.78, y, 11.75, 0.54, fill)
        add_text(slide, title, 1.05, y + 0.11, 1.55, 0.18, 8.7, "deep" if fill != "white" else "green", True)
        add_text(slide, body, 2.72, y + 0.16, 8.95, 0.12, 8.8, "deep" if fill != "white" else "muted")
        if i < len(layers) - 1:
            add_arrow(slide, 6.65, y + 0.55, 6.65, y + 0.73, "pale_green", 1.4)

    add_round(slide, 0.9, 6.15, 11.45, 0.52, "deep")
    add_text(slide, "关键难点：前端交互状态、Agent stage 追踪、多模态结构化、推荐 fallback、确认摄入边界、长期营养/偏好记忆一致性。", 1.22, 6.34, 10.8, 0.12, 9.6, "white", True, PP_ALIGN.CENTER)


def slide_agent_deep(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, "bg")
    add_title(slide, "Agent 深度：不是聊天框，而是可追踪的工具编排", "独立角色页面保证可测；Agent 负责跨场景理解、调用工具、生成结果和学习反馈")

    add_round(slide, 0.75, 1.45, 3.35, 4.6, "white")
    add_text(slide, "1. 输入理解", 1.05, 1.78, 1.5, 0.25, 15, "green", True)
    add_text(slide, "自然语言 + 图片输入\nDeepSeek -> Ollama -> 规则兜底\n输出结构化意图：食材、目标、时间、模式", 1.05, 2.22, 2.6, 0.78, 11.5, "muted")
    add_round(slide, 1.05, 3.25, 2.48, 0.58, "pale_green")
    add_text(slide, "例：牛肉南瓜，30分钟减脂餐", 1.22, 3.48, 2.1, 0.12, 9.6, "deep", True, PP_ALIGN.CENTER)
    add_text(slide, "价值：用户不需要理解系统有多少页面，Agent 先把任务拆开。", 1.05, 4.18, 2.6, 0.45, 10.5, "muted")

    add_round(slide, 4.42, 1.45, 4.15, 4.6, "deep")
    add_text(slide, "2. 工具编排", 4.78, 1.78, 1.5, 0.25, 15, "white", True)
    stages = [
        ("sense", "VLM 识别食材/菜品/品质"),
        ("decision", "食材 + 目标 + 偏好 + 缺口推荐"),
        ("task", "生成清单、计划、库存扣减"),
        ("feedback", "评分文本 -> 偏好记忆"),
    ]
    for i, (name, desc) in enumerate(stages):
        y = 2.25 + i * 0.68
        add_round(slide, 4.85, y, 3.18, 0.46, "white")
        add_text(slide, name, 5.05, y + 0.16, 0.72, 0.1, 8.8, "green", True)
        add_text(slide, desc, 5.78, y + 0.16, 2.02, 0.1, 8.6, "deep")
        if i < len(stages) - 1:
            add_arrow(slide, 6.42, y + 0.47, 6.42, y + 0.66, "pale_green", 1.4)
    add_text(slide, "每个 stage 返回 status / latency_ms / degraded，让演示时能解释哪里慢、哪里降级。", 4.85, 5.35, 3.15, 0.32, 9.5, RGBColor(217, 239, 229))

    add_round(slide, 8.9, 1.45, 3.7, 4.6, "white")
    add_text(slide, "3. 长期记忆回流", 9.22, 1.78, 1.8, 0.25, 15, "green", True)
    add_text(slide, "确认摄入写入 MealRecord\n每日/每周营养看板读取 completed 记录\n评分原因写入 preference_memories", 9.22, 2.23, 2.8, 0.68, 11.2, "muted")
    add_round(slide, 9.22, 3.25, 2.78, 0.86, "pale_amber")
    add_text(slide, "喜欢清淡少油、高蛋白\n=> liked: light, high_protein", 9.45, 3.48, 2.3, 0.24, 9.5, "deep", True, PP_ALIGN.CENTER)
    add_round(slide, 9.22, 4.42, 2.78, 0.86, "pale_purple")
    add_text(slide, "太油腻、不想再吃\n=> avoid: oily", 9.45, 4.68, 2.3, 0.2, 9.5, "deep", True, PP_ALIGN.CENTER)
    add_text(slide, "下一轮推荐读取 liked/avoid，对菜谱加权或降权。", 9.22, 5.48, 2.8, 0.22, 10, "muted")


def slide_tests_demo_scale(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, "bg")
    add_title(slide, "角色化测试、软件规模与现场演示", "测试跟功能一一呼应，最后用一条主线演示 B-Y-T-E 闭环")

    flow = [
        ("1", "画像", "目标/身高体重"),
        ("2", "识别", "删改去重"),
        ("3", "清单", "营养占比"),
        ("4", "摄入", "写入看板"),
        ("5", "反馈", "偏好记忆"),
    ]
    for i, (n, title, body) in enumerate(flow):
        x = 0.72 + i * 2.48
        add_round(slide, x, 1.45, 2.08, 0.92, "white")
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.18), Inches(1.67), Inches(0.34), Inches(0.34))
        circ.fill.solid(); circ.fill.fore_color.rgb = COLORS["green"]; circ.line.fill.background()
        add_text(slide, n, x + 0.27, 1.76, 0.16, 0.1, 7.8, "white", True, PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.62, 1.67, 0.85, 0.18, 11, "deep", True)
        add_text(slide, body, x + 0.62, 1.96, 1.05, 0.14, 8.4, "muted")
        if i < len(flow) - 1:
            add_arrow(slide, x + 2.1, 1.92, x + 2.42, 1.92, "green", 1.8)

    rows = [
        ("家庭做饭", "场景一", "识别/校正/推荐"),
        ("健康管理", "场景四", "确认摄入/删除"),
        ("买菜用户", "场景三", "品质鉴定"),
        ("探店用户", "场景五", "故事口味吃法"),
        ("长期个性化", "pytest", "偏好记忆"),
    ]
    add_round(slide, 0.82, 2.85, 6.6, 2.18, "white")
    add_text(slide, "测试与功能对应", 1.1, 3.08, 1.6, 0.2, 13.5, "deep", True)
    for i, row in enumerate(rows):
        y = 3.45 + i * 0.28
        add_text(slide, row[0], 1.12, y, 1.1, 0.09, 7.5, "deep", True)
        add_text(slide, row[1], 2.55, y, 0.65, 0.09, 7.4, "green", True)
        add_text(slide, row[2], 3.45, y, 1.65, 0.09, 7.4, "muted")

    metrics = [("测试文件", "11", ">3"), ("核心源文件", "98", ">5"), ("核心代码", "12,295", ">500")]
    for i, (label, value, req) in enumerate(metrics):
        x = 7.78 + i * 1.48
        add_round(slide, x, 2.85, 1.28, 1.06, "deep" if i == 0 else "white")
        tc = "white" if i == 0 else "deep"
        add_text(slide, label, x + 0.12, 3.1, 1.02, 0.1, 7.8, tc, True, PP_ALIGN.CENTER)
        add_text(slide, value, x + 0.1, 3.37, 1.08, 0.18, 12.8 if i < 2 else 10.5, tc, True, PP_ALIGN.CENTER)
        add_text(slide, f"要求 {req}", x + 0.1, 3.64, 1.08, 0.08, 6.8, tc if i == 0 else "muted", False, PP_ALIGN.CENTER)
    add_round(slide, 7.78, 4.28, 4.24, 0.75, "pale_green")
    add_text(slide, "分布：后端 3,636 行 / 前端 7,935 行 / 测试 724 行；菜谱数据 138,252 行。", 7.98, 4.48, 3.85, 0.16, 8.2, "deep", True, PP_ALIGN.CENTER)
    add_text(slide, "最新验证：15 个相关 pytest 通过；两个 H5 前端构建通过。", 8.05, 4.76, 3.72, 0.12, 8.4, "muted", False, PP_ALIGN.CENTER)

    add_round(slide, 1.02, 5.58, 11.22, 0.72, "deep")
    add_text(slide, "收束：ByteSavor 不是只接几个 AI 接口，而是把饮食任务做成“感知 -> 决策 -> 执行 -> 反馈”的可验证闭环。", 1.34, 5.81, 10.6, 0.18, 11.7, "white", True, PP_ALIGN.CENTER)
    add_text(slide, "兜底：若现场模型延迟，切换 demo_tests 图片、Swagger 接口或 pytest/build 验证结果。", 1.35, 6.55, 10.5, 0.2, 10.8, "muted", False, PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide_cover(prs)
    slide_baseline_to_closure(prs)
    slide_arch(prs)
    slide_agent_deep(prs)
    slide_tests_demo_scale(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


OUT = Path("/Users/liwenbin930/Desktop/bytesavor-backend/期末展演/ByteSavor_可编辑图表资产.pptx")

W, H = 13.333, 7.5

COLORS = {
    "bg": RGBColor(247, 251, 248),
    "deep": RGBColor(23, 59, 46),
    "green": RGBColor(27, 127, 95),
    "mint": RGBColor(85, 201, 155),
    "amber": RGBColor(244, 185, 87),
    "purple": RGBColor(92, 86, 177),
    "muted": RGBColor(92, 107, 99),
    "pale_green": RGBColor(234, 247, 239),
    "pale_amber": RGBColor(255, 248, 234),
    "pale_purple": RGBColor(244, 243, 255),
    "white": RGBColor(255, 255, 255),
}


def add_text(slide, text, x, y, w, h, size=18, color="deep", bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "PingFang SC"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = COLORS[color] if isinstance(color, str) else color
    if align:
        p.alignment = align
    return box


def add_round(slide, x, y, w, h, fill="white", line=None, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill] if isinstance(fill, str) else fill
    if line:
        shape.line.color.rgb = COLORS[line] if isinstance(line, str) else line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_title(slide, title, subtitle):
    add_text(slide, title, 0.75, 0.36, 11.8, 0.44, 25, "deep", True)
    add_text(slide, subtitle, 0.76, 0.82, 11.5, 0.32, 13.5, "muted")


def add_arrow(slide, x1, y1, x2, y2, color="muted", width=2.2):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = COLORS[color]
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def slide_baseline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["bg"]
    add_title(slide, "从单点工具到全链路 Agent", "ByteSavor 把识别、推荐、执行和反馈接成一条可验证流程")

    add_text(slide, "传统饮食软件：能力割裂", 1.05, 1.55, 4.2, 0.35, 17, "muted", True)
    add_text(slide, "ByteSavor：场景闭环", 7.45, 1.55, 4.2, 0.35, 17, "deep", True)

    cards = [
        ("热量记录", "只能手动估算", 0.95, 2.18),
        ("菜谱搜索", "不知道家里有什么", 3.45, 2.18),
        ("购物清单", "和推荐断开", 0.95, 3.75),
        ("图片识别", "识别后没有执行链路", 3.45, 3.75),
    ]
    for title, sub, x, y in cards:
        add_round(slide, x, y, 2.15, 1.04)
        add_text(slide, title, x + 0.23, y + 0.2, 1.7, 0.28, 15, "deep", True)
        add_text(slide, sub, x + 0.23, y + 0.55, 1.7, 0.24, 11, "muted")
    add_arrow(slide, 3.12, 2.7, 3.4, 2.7, "muted", 1.5)
    add_arrow(slide, 2.02, 3.24, 2.02, 3.72, "muted", 1.5)
    add_arrow(slide, 4.52, 3.24, 4.52, 3.72, "muted", 1.5)
    add_arrow(slide, 3.12, 4.27, 3.4, 4.27, "muted", 1.5)

    add_round(slide, 7.15, 2.05, 5.25, 3.75)
    cx, cy = 9.78, 3.92
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 0.84), Inches(cy - 0.84), Inches(1.68), Inches(1.68))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS["green"]
    circle.line.fill.background()
    add_text(slide, "B-Y-T-E", cx - 0.65, cy - 0.18, 1.3, 0.28, 18, "white", True, PP_ALIGN.CENTER)
    add_text(slide, "饮食闭环", cx - 0.65, cy + 0.16, 1.3, 0.22, 11, "white", False, PP_ALIGN.CENTER)

    nodes = [
        ("B 感知", cx - 0.52, cy - 1.6, "green"),
        ("Y 决策", cx + 1.65, cy - 0.05, "amber"),
        ("T 执行", cx - 0.52, cy + 1.52, "purple"),
        ("E 反馈", cx - 2.12, cy - 0.05, "green"),
    ]
    for label, x, y, c in nodes:
        add_round(slide, x, y, 1.15, 0.38, "pale_green" if c == "green" else "pale_amber" if c == "amber" else "pale_purple")
        add_text(slide, label, x + 0.05, y + 0.07, 1.05, 0.22, 11.5, c, True, PP_ALIGN.CENTER)
    add_arrow(slide, cx, cy - 1.12, cx + 1.52, cy - 0.18, "green", 2)
    add_arrow(slide, cx + 1.52, cy + 0.18, cx, cy + 1.2, "amber", 2)
    add_arrow(slide, cx - 0.25, cy + 1.2, cx - 1.55, cy + 0.18, "purple", 2)
    add_arrow(slide, cx - 1.55, cy - 0.18, cx - 0.25, cy - 1.12, "green", 2)
    add_text(slide, "确认摄入、评分原因、偏好记忆会回到下一轮推荐", 7.85, 6.02, 4.1, 0.35, 12, "muted", False, PP_ALIGN.CENTER)


def slide_arch(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["deep"]
    add_text(slide, "B-Y-T-E Agent 技术闭环", 0.75, 0.36, 11.8, 0.44, 25, "white", True)
    add_text(slide, "每个工具独立可测，由 Agent 统一编排，并把反馈写回长期记忆", 0.76, 0.82, 11.5, 0.32, 13.5, RGBColor(201, 225, 213))
    add_round(slide, 4.5, 1.45, 4.35, 0.65, "white")
    add_text(slide, "uni-app / Vue3 前端", 4.8, 1.68, 3.75, 0.25, 16, "deep", True, PP_ALIGN.CENTER)
    add_arrow(slide, 6.67, 2.1, 6.67, 2.75, "pale_green", 2.5)
    modules = [
        ("B 感知", "图片/文本输入\n食材、分量、品质、菜品\nVLM: Qwen-VL", 0.9, "green"),
        ("Y 决策", "推荐引擎\n目标、偏好、营养缺口\n食材50% + 标签30% + 偏好20%", 3.55, "amber"),
        ("T 执行", "三餐计划 / 清单\n库存扣减、完成记录\n识别不等于摄入", 6.2, "purple"),
        ("E 反馈", "评分 + 文字原因\nLLM 解析偏好记忆\nliked / avoid signals", 8.85, "green"),
    ]
    for title, body, x, c in modules:
        add_round(slide, x, 2.9, 2.35, 1.9, "white")
        add_text(slide, title, x + 0.22, 3.16, 1.8, 0.3, 18, c, True)
        add_text(slide, body, x + 0.22, 3.58, 1.9, 0.85, 10.8, "muted")
    add_arrow(slide, 3.25, 3.85, 3.52, 3.85, "pale_green", 2)
    add_arrow(slide, 5.9, 3.85, 6.17, 3.85, "pale_green", 2)
    add_arrow(slide, 8.55, 3.85, 8.82, 3.85, "pale_green", 2)
    add_arrow(slide, 10.02, 4.82, 2.05, 5.9, "pale_green", 2)
    add_text(slide, "反馈回路：确认摄入 -> 营养长期记录；评分原因 -> 偏好数据库 -> 下一轮推荐加权", 2.1, 5.96, 9.2, 0.3, 12, RGBColor(217, 239, 229), False, PP_ALIGN.CENTER)
    infra = ["FastAPI", "MySQL", "Redis", "VLM / LLM", "2576+ 菜谱"]
    for i, item in enumerate(infra):
        x = 1.25 + i * 2.15
        add_round(slide, x, 6.55, 1.55, 0.45, "white")
        add_text(slide, item, x + 0.05, 6.68, 1.45, 0.18, 11.5, "deep", True, PP_ALIGN.CENTER)


def slide_matrix(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["bg"]
    add_title(slide, "角色化功能测试矩阵", "功能展示和测试场景一一对应，现场演示可以按角色切换")
    headers = ["角色", "功能路径", "测试材料", "验证点"]
    widths = [1.65, 3.65, 2.65, 3.6]
    xs = [0.95, 2.6, 6.25, 8.9]
    add_round(slide, 0.85, 1.65, 11.65, 0.48, "deep")
    for h, x, w in zip(headers, xs, widths):
        add_text(slide, h, x, 1.8, w, 0.18, 11, "white", True)
    rows = [
        ("家庭做饭", "拍照识别 -> 校正 -> 推荐 -> 清单", "场景一_拍照推荐", "可删改、推荐多道菜、不误计入"),
        ("健康管理", "餐盘识别 -> 确认摄入 -> 看板", "场景四_营养分析", "确认才写入；错误可删除"),
        ("买菜用户", "拍水果/食材 -> 品质鉴定", "场景三_品质鉴定", "等级、依据、购买建议"),
        ("探店用户", "拍菜品 -> 故事/口味/吃法", "场景五_探店向导", "VLM 识别 + LLM 补全讲解"),
        ("备餐采购", "多菜谱/文本 -> 合并购物清单", "场景二_清单导出", "同名合并、数量保留、可复制"),
        ("长期个性化", "确认摄入 -> 评分原因 -> 偏好记忆", "test_feedback_memory.py", "liked/avoid 写库并影响推荐"),
    ]
    fills = ["pale_green", "pale_amber", "pale_purple", "pale_green", "pale_amber", "pale_green"]
    for i, row in enumerate(rows):
        y = 2.3 + i * 0.74
        add_round(slide, 0.85, y, 11.65, 0.56, fills[i])
        add_text(slide, row[0], xs[0], y + 0.18, widths[0], 0.16, 10.8, "deep", True)
        add_text(slide, row[1], xs[1], y + 0.18, widths[1], 0.16, 10.3, "deep")
        add_text(slide, row[2], xs[2], y + 0.18, widths[2], 0.16, 10.3, "green" if i in (0, 3, 5) else "amber" if i in (1, 4) else "purple", True)
        add_text(slide, row[3], xs[3], y + 0.18, widths[3], 0.16, 10.3, "deep")


def slide_scale(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["bg"]
    add_title(slide, "软件规模与验证结果", "规模满足课程要求，关键流程有自动化测试和构建验证")
    metrics = [
        ("自动化测试文件", "11", "个", "覆盖 auth / decision / meals / feedback", "deep", "white"),
        ("主要源文件", "105", "个", "后端 app + 前端 bsapp/src", "white", "deep"),
        ("代码总行数", "149,823", "", "远超 500 行要求", "white", "deep"),
    ]
    for i, (label, num, unit, foot, fill, text_color) in enumerate(metrics):
        x = 0.9 + i * 4.05
        add_round(slide, x, 1.8, 3.3, 1.55, fill)
        add_text(slide, label, x + 0.35, 2.15, 2.5, 0.22, 12.5, text_color)
        add_text(slide, num, x + 0.35, 2.48, 2.35, 0.45, 33, text_color, True)
        if unit:
            add_text(slide, unit, x + 1.55, 2.73, 0.45, 0.2, 13, text_color)
        add_text(slide, foot, x + 0.35, 3.05, 2.6, 0.22, 10.5, text_color if fill == "deep" else "muted")
    add_round(slide, 0.9, 4.05, 11.55, 2.1, "white")
    add_text(slide, "关键验证链路", 1.25, 4.42, 3.2, 0.28, 16, "deep", True)
    items = [
        ("1", "个性化目标", "身高体重、运动频次、自定义目标", "green", "pale_green"),
        ("2", "确认摄入", "只有确认后才写入长期营养", "amber", "pale_amber"),
        ("3", "偏好记忆", "评分文字解析为 liked / avoid", "purple", "pale_purple"),
        ("4", "前端构建", "两个 uni-app 目录均 build 通过", "green", "pale_green"),
    ]
    for i, (n, title, sub, c, fill) in enumerate(items):
        x = 1.25 + i * 2.85
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(4.92), Inches(0.38), Inches(0.38))
        circ.fill.solid()
        circ.fill.fore_color.rgb = COLORS[fill]
        circ.line.fill.background()
        add_text(slide, n, x + 0.02, 5.02, 0.34, 0.12, 9.5, c, True, PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.5, 4.9, 1.7, 0.23, 12, "deep", True)
        add_text(slide, sub, x + 0.5, 5.22, 2.05, 0.35, 9.5, "muted")
    add_text(slide, "最新验证：15 个相关 pytest 通过；H5 主前端和嵌套前端均 DONE Build complete。", 1.25, 5.82, 9.6, 0.22, 11, "muted")


def main():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide_baseline(prs)
    slide_arch(prs)
    slide_matrix(prs)
    slide_scale(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
